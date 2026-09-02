import os, json, re, uuid, io
from datetime import datetime
from fastapi import FastAPI, Request
from fastapi.responses import JSONResponse, StreamingResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
import sqlalchemy as sa
from sqlalchemy import text

GROQ_API_KEY = os.environ.get('GROQ_API_KEY', '')
DATABASE_URL = os.environ.get('DATABASE_URL', 'sqlite:///./memory.db')

client = Groq(api_key=GROQ_API_KEY) if GROQ_API_KEY else None
engine = sa.create_engine(DATABASE_URL, pool_pre_ping=True)

app = FastAPI(title='Jarvis Assistant')
app.add_middleware(CORSMiddleware, allow_origins=['*'], allow_methods=['*'], allow_headers=['*'])

os.makedirs('generated', exist_ok=True)

def init_db():
    with engine.begin() as conn:
        if DATABASE_URL.startswith('sqlite'):
            conn.execute(text('''CREATE TABLE IF NOT EXISTS messages (id INTEGER PRIMARY KEY AUTOINCREMENT, session_id TEXT NOT NULL, role TEXT NOT NULL, content TEXT NOT NULL, created_at TEXT DEFAULT CURRENT_TIMESTAMP)'''))
            conn.execute(text('''CREATE TABLE IF NOT EXISTS facts (id INTEGER PRIMARY KEY AUTOINCREMENT, session_id TEXT NOT NULL, fact TEXT NOT NULL)'''))
        else:
            conn.execute(text('''CREATE TABLE IF NOT EXISTS messages (id SERIAL PRIMARY KEY, session_id TEXT NOT NULL, role TEXT NOT NULL, content TEXT NOT NULL, created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP)'''))
            conn.execute(text('''CREATE TABLE IF NOT EXISTS facts (id SERIAL PRIMARY KEY, session_id TEXT NOT NULL, fact TEXT NOT NULL)'''))

init_db()

def save_message(session_id, role, content):
    with engine.begin() as conn:
        conn.execute(text('INSERT INTO messages (session_id, role, content) VALUES (:s,:r,:c)'), {'s': session_id, 'r': role, 'c': content})

def get_history(session_id, limit=40):
    with engine.begin() as conn:
        rows = conn.execute(text('SELECT role, content FROM messages WHERE session_id=:s ORDER BY id DESC LIMIT :l'), {'s': session_id, 'l': limit}).fetchall()
    return [{'role': r[0], 'content': r[1]} for r in reversed(rows)]

def get_facts(session_id):
    with engine.begin() as conn:
        rows = conn.execute(text('SELECT fact FROM facts WHERE session_id=:s'), {'s': session_id}).fetchall()
    return [r[0] for r in rows]

def add_fact(session_id, fact):
    with engine.begin() as conn:
        conn.execute(text('INSERT INTO facts (session_id, fact) VALUES (:s,:f)'), {'s': session_id, 'f': fact})

SYSTEM_PROMPT = """Tu es Jarvis, l'assistant IA personnel de l'utilisateur. Tu as une memoire persistante: utilise les faits connus fournis pour personnaliser tes reponses. Reponds en francais, de facon directe et utile.

Si l'utilisateur demande du CODE (jeu, script, page web, appli simple): reponds avec une courte explication puis un bloc de code HTML/CSS/JS AUTONOME et complet dans un seul bloc ```html ... ``` qui peut tourner seul dans un iframe (pas de dependances externes sauf CDN standard). Le jeu/appli doit etre jouable directement.

Si l'utilisateur demande une PRESENTATION/PowerPoint/expose: reponds avec une courte phrase puis un bloc JSON dans ```json ... ``` avec ce format exact: {"slides": [{"title": "...", "bullets": ["...", "..."]}]}

Si l'utilisateur demande une IMAGE: reponds avec une courte phrase puis un bloc dans ```image ... ``` contenant juste une description en anglais precise pour un generateur d'images.

Si l'utilisateur te donne une information personnelle importante a retenir durablement, ajoute a la fin de ta reponse une ligne exacte: [FACT: le fait resume en une phrase]
"""

class ChatIn(BaseModel):
    session_id: str
    message: str

def extract_block(text_, lang):
    m = re.search(rf'```{lang}\s*(.*?)```', text_, re.DOTALL)
    return m.group(1).strip() if m else None

def extract_fact(text_):
    m = re.search(r'\[FACT:\s*(.*?)\]', text_)
    return m.group(1).strip() if m else None

@app.post('/api/chat')
def chat(payload: ChatIn):
    if client is None:
        return JSONResponse({'error': 'GROQ_API_KEY manquante sur le serveur'}, status_code=500)

    session_id = payload.session_id
    save_message(session_id, 'user', payload.message)

    facts = get_facts(session_id)
    facts_block = '\n'.join(f'- {f}' for f in facts) if facts else '(aucun fait connu encore)'
    history = get_history(session_id)

    messages = [{'role': 'system', 'content': SYSTEM_PROMPT + '\n\nFaits connus sur l\'utilisateur:\n' + facts_block}]
    messages += history

    completion = client.chat.completions.create(model='llama-3.3-70b-versatile', messages=messages, temperature=0.7, max_tokens=4000)
    reply = completion.choices[0].message.content

    save_message(session_id, 'assistant', reply)

    fact = extract_fact(reply)
    if fact:
        add_fact(session_id, fact)
        reply = reply.replace(f'[FACT: {fact}]', '').strip()

    artifact = None
    code_block = extract_block(reply, 'html')
    json_block = extract_block(reply, 'json')
    image_block = extract_block(reply, 'image')

    if code_block:
        artifact = {'type': 'code', 'content': code_block}
        reply = re.sub(r'```html.*?```', '', reply, flags=re.DOTALL).strip()
    elif json_block:
        try:
            slides = json.loads(json_block)
            artifact = {'type': 'slides', 'content': slides}
        except Exception:
            pass
        reply = re.sub(r'```json.*?```', '', reply, flags=re.DOTALL).strip()
    elif image_block:
        artifact = {'type': 'image', 'content': image_block.strip()}
        reply = re.sub(r'```image.*?```', '', reply, flags=re.DOTALL).strip()

    return {'reply': reply, 'artifact': artifact}


@app.get('/api/image')
def image_proxy(prompt: str):
    import urllib.parse, urllib.request
    url = 'https://image.pollinations.ai/prompt/' + urllib.parse.quote(prompt) + '?width=768&height=768&nologo=true'
    req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
    data = urllib.request.urlopen(req, timeout=60).read()
    return StreamingResponse(io.BytesIO(data), media_type='image/jpeg')


@app.post('/api/pptx')
def build_pptx(payload: dict):
    slides_data = payload.get('slides', [])
    prs = Presentation()
    for s in slides_data:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = s.get('title', '')
        body = slide.placeholders[1].text_frame
        bullets = s.get('bullets', [])
        for i, b in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = b
            p.font.size = Pt(20)

    file_id = str(uuid.uuid4())
    path = f'generated/{file_id}.pptx'
    prs.save(path)
    return {'download_url': f'/api/download/{file_id}.pptx'}


@app.get('/api/download/{filename}')
def download(filename: str):
    path = f'generated/{filename}'
    return FileResponse(path, filename=filename)


app.mount('/', StaticFiles(directory='static', html=True), name='static')
